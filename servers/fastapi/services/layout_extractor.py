import os
import uuid
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


EMU_PER_INCH = 914400
DEFAULT_DPI = 96


def emu_to_px(emu: int, dpi: int = DEFAULT_DPI) -> int:
    """Convert EMU to pixels."""
    return int(round(emu * dpi / EMU_PER_INCH))


def _color_to_hex(color) -> Optional[str]:
    """Convert pptx color to #RRGGBB when possible."""
    if color is None:
        return None
    try:
        if color.type == MSO_COLOR_TYPE.RGB and color.rgb:
            rgb = color.rgb
            return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return None
    return None


def _font_spec(font) -> Dict[str, Any]:
    family = font.name
    size = round(float(font.size.pt), 2) if font.size else None
    weight = 700 if font.bold else 400
    style = "italic" if font.italic else "normal"
    underline = bool(font.underline) if font.underline is not None else False
    color = _color_to_hex(font.color)
    return {
        "family": family,
        "size": size,
        "weight": weight,
        "style": style,
        "underline": underline,
        "color": color,
    }


def _paragraph_marker(paragraph) -> str:
    """Detect bullet/number/none."""
    ppr = getattr(paragraph, "_pPr", None)
    if ppr is None:
        return "bullet" if paragraph.level else "none"

    if ppr.find(qn("a:buNone")) is not None:
        return "none"
    if ppr.find(qn("a:buAutoNum")) is not None:
        return "number"
    if ppr.find(qn("a:buChar")) is not None:
        return "bullet"
    return "bullet" if paragraph.level else "none"


def _paragraph_runs(paragraph) -> List[Dict[str, Any]]:
    runs = []
    for run in paragraph.runs:
        text = run.text or ""
        if text == "":
            continue
        font = _font_spec(run.font)
        run_color = _color_to_hex(run.font.color) or font.get("color")
        runs.append(
            {
                "text": text,
                "font": font,
                "color": run_color,
                "highlight": None,
                "underline": bool(run.font.underline) if run.font.underline is not None else False,
                "strike": bool(getattr(run.font, "strike", False)),
            }
        )
    return runs


def _text_element_from_shape(shape, z: int):
    bbox = {
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "width": emu_to_px(shape.width),
        "height": emu_to_px(shape.height),
    }

    paragraph_blocks = []
    for paragraph in shape.text_frame.paragraphs:
        marker = _paragraph_marker(paragraph)
        runs = _paragraph_runs(paragraph)
        if not runs and not paragraph.text:
            continue
        if not runs and paragraph.text:
            font = _font_spec(paragraph.font)
            runs.append(
                {
                    "text": paragraph.text,
                    "font": font,
                    "color": _color_to_hex(paragraph.font.color) or font.get("color"),
                    "highlight": None,
                    "underline": bool(paragraph.font.underline) if paragraph.font.underline is not None else False,
                    "strike": bool(getattr(paragraph.font, "strike", False)),
                }
            )

        paragraph_blocks.append(
            {
                "level": paragraph.level or 0,
                "marker": marker,
                "text": paragraph.text or "",
                "runs": runs,
            }
        )

    alignment = getattr(shape.text_frame, "paragraphs", [None])[0]
    align_value = getattr(alignment, "alignment", None)
    align_map = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
        PP_ALIGN.DISTRIBUTE: "justify",
    }
    align = align_map.get(align_value, "left")

    return {
        "id": str(uuid.uuid4()),
        "type": "text",
        "bbox": bbox,
        "z": z,
        "rotation": getattr(shape, "rotation", None) or 0,
        "opacity": 1,
        "align": align,
        "runs": [run for block in paragraph_blocks for run in block.get("runs", [])],
        "bullets": paragraph_blocks,
    }


def _fill_from_shape(shape) -> Optional[Dict[str, Any]]:
    """Extract fill information from a shape, handling all fill types safely."""

    # FIRST: Try direct XML access (most reliable for shapes created in PowerPoint)
    try:
        sp_pr = shape._element.spPr
        if sp_pr is not None:
            # Look for solidFill in XML
            solid_fill = sp_pr.find(qn('a:solidFill'))
            if solid_fill is not None:
                # Try sRGB color
                srgb_clr = solid_fill.find(qn('a:srgbClr'))
                if srgb_clr is not None and 'val' in srgb_clr.attrib:
                    hex_color = '#' + srgb_clr.attrib['val'].upper()
                    return {"type": "solid", "color": hex_color}

                # Try scheme color (theme color) - resolve via python-pptx
                scheme_clr = solid_fill.find(qn('a:schemeClr'))
                if scheme_clr is not None and 'val' in scheme_clr.attrib:
                    # Get the theme color name (e.g., "accent4")
                    theme_color_name = scheme_clr.attrib['val']

                    # Try to resolve via python-pptx API
                    try:
                        fill = shape.fill
                        fill_type = fill.type
                        if fill_type == MSO_FILL.SOLID:
                            color = _color_to_hex(fill.fore_color)
                            if color:
                                return {"type": "solid", "color": color}
                    except:
                        pass
    except:
        pass

    # SECOND: Try python-pptx high-level API
    try:
        fill = shape.fill
    except Exception:
        return None

    if fill is None:
        return None

    try:
        # Get fill type
        fill_type = fill.type

        if fill_type is None:
            return None

        # Solid fill
        if fill_type == MSO_FILL.SOLID:
            try:
                color = _color_to_hex(fill.fore_color)
                if color:
                    return {"type": "solid", "color": color}
            except (TypeError, AttributeError):
                return None

        # Gradient fill
        elif fill_type == MSO_FILL.GRADIENT:
            try:
                color = _color_to_hex(fill.fore_color)
                if color:
                    return {"type": "solid", "color": color}
            except (TypeError, AttributeError):
                return None

        # Pattern fill
        elif fill_type == MSO_FILL.PATTERNED:
            try:
                color = _color_to_hex(fill.fore_color)
                if color:
                    return {"type": "solid", "color": color}
            except (TypeError, AttributeError):
                return None

        # Background fill (inherits from theme/master)
        elif fill_type == MSO_FILL.BACKGROUND:
            # Try to get the effective color by accessing underlying XML
            try:
                # First try: use .solid() to convert to solid and get color
                fill.solid()
                color = _color_to_hex(fill.fore_color)
                if color:
                    return {"type": "solid", "color": color}
            except:
                pass

            # Second try: access the shape's spPr (shape properties) XML directly
            try:
                sp_pr = shape._element.spPr
                if sp_pr is not None:
                    # Look for solidFill in XML
                    solid_fill = sp_pr.find(qn('a:solidFill'))
                    if solid_fill is not None:
                        # Get color from XML
                        srgb_clr = solid_fill.find(qn('a:srgbClr'))
                        if srgb_clr is not None and 'val' in srgb_clr.attrib:
                            hex_color = '#' + srgb_clr.attrib['val']
                            return {"type": "solid", "color": hex_color}
            except:
                pass

            return None

        # Picture/texture fill
        elif fill_type == MSO_FILL.PICTURE:
            return None

    except Exception:
        pass

    return None


def _stroke_from_shape(shape) -> Optional[Dict[str, Any]]:
    try:
        line = shape.line
    except Exception:
        return None
    if line is None:
        return None
    width = line.width.pt if line.width else None
    color = _color_to_hex(getattr(line, "color", None))
    if width is None and color is None:
        return None
    return {"width": width, "color": color}


def _image_element(shape, asset_dir: str, asset_url_prefix: str, slide_index: int, z: int):
    img = shape.image
    ext = img.ext or "png"
    filename = f"slide_{slide_index}_image_{z}.{ext}"
    path = os.path.join(asset_dir, filename)
    with open(path, "wb") as fp:
        fp.write(img.blob)

    bbox = {
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "width": emu_to_px(shape.width),
        "height": emu_to_px(shape.height),
    }

    return {
        "id": str(uuid.uuid4()),
        "type": "image",
        "bbox": bbox,
        "z": z,
        "rotation": getattr(shape, "rotation", None) or 0,
        "opacity": 1,
        "src": f"{asset_url_prefix}/{filename}",
        "object_fit": "cover",
    }


def _table_element(shape, z: int):
    bbox = {
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "width": emu_to_px(shape.width),
        "height": emu_to_px(shape.height),
    }
    rows_data = []
    table = shape.table
    for row in table.rows:
        cells = []
        for cell in row.cells:
            runs = []
            for paragraph in cell.text_frame.paragraphs:
                runs.extend(_paragraph_runs(paragraph))
            cells.append({"text": cell.text, "runs": runs})
        rows_data.append({"cells": cells})

    return {
        "id": str(uuid.uuid4()),
        "type": "table",
        "bbox": bbox,
        "z": z,
        "rows": rows_data,
    }


def _shape_element(shape, z: int):
    bbox = {
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "width": emu_to_px(shape.width),
        "height": emu_to_px(shape.height),
    }
    fill = _fill_from_shape(shape)
    stroke = _stroke_from_shape(shape)
    return {
        "id": str(uuid.uuid4()),
        "type": "shape",
        "shape": getattr(shape, "shape_type", None).name
        if getattr(shape, "shape_type", None)
        else "shape",
        "bbox": bbox,
        "z": z,
        "rotation": getattr(shape, "rotation", None) or 0,
        "opacity": 1,
        "fill": fill,
        "stroke": stroke,
    }


@dataclass
class ParsedSlide:
    width_px: int
    height_px: int
    background: Optional[Dict[str, Any]]
    fonts: List[str] = field(default_factory=list)
    elements: List[Dict[str, Any]] = field(default_factory=list)
    index: int = 0
    id: str = field(default_factory=lambda: str(uuid.uuid4()))


def parse_pptx_to_layouts(
    pptx_path: str, asset_output_dir: str, asset_url_prefix: str
) -> List[ParsedSlide]:
    """
    Parse a PPTX into slide layouts with deterministic geometry and extracted assets.
    """
    presentation = Presentation(pptx_path)
    slide_width_px = emu_to_px(presentation.slide_width)
    slide_height_px = emu_to_px(presentation.slide_height)

    os.makedirs(asset_output_dir, exist_ok=True)

    slides: List[ParsedSlide] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        elements: List[Dict[str, Any]] = []
        fonts: List[str] = []

        background_fill = getattr(slide, "background", None)
        bg_color = None
        if background_fill and getattr(background_fill, "fill", None):
            try:
                bg_color = _color_to_hex(background_fill.fill.fore_color)
            except (TypeError, AttributeError):
                # Fill has no foreground color (NoFill, or other fill types)
                bg_color = None
        background = {"type": "solid", "color": bg_color} if bg_color else None

        for z, shape in enumerate(slide.shapes):
            # Debug logging
            import logging
            logger = logging.getLogger(__name__)

            # Skip placeholder shapes with no content
            is_placeholder = getattr(shape, "is_placeholder", False)

            logger.info(f"Shape {z}: type={shape.shape_type}, has_text_frame={shape.has_text_frame if hasattr(shape, 'has_text_frame') else False}, is_placeholder={is_placeholder}")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                elements.append(
                    _image_element(shape, asset_output_dir, asset_url_prefix, slide_index, z)
                )
                continue

            if shape.has_table:
                elements.append(_table_element(shape, z))
                continue

            if shape.has_text_frame:
                # Shape with text - extract BOTH fill/stroke AND text content
                text_el = _text_element_from_shape(shape, z)

                # Check if this shape has actual content (not just empty placeholder)
                has_text = any(run.get("text", "").strip() for run in text_el.get("runs", []))

                # Also extract shape fill and stroke
                fill = _fill_from_shape(shape)
                stroke = _stroke_from_shape(shape)

                # Skip empty placeholders: placeholder with no text, no fill, no stroke
                if is_placeholder and not has_text and not fill and not stroke:
                    logger.info(f"  → Skipping empty placeholder shape {z}")
                    continue

                # Add fill and stroke to the text element
                if fill:
                    text_el["fill"] = fill
                if stroke:
                    text_el["stroke"] = stroke

                # Only add if: has text content OR has fill/stroke (visible shape)
                if has_text or fill or stroke:
                    for run in text_el.get("runs", []):
                        family = run.get("font", {}).get("family")
                        if family:
                            fonts.append(family)
                    elements.append(text_el)
                else:
                    logger.info(f"  → Skipping shape {z} with no visible content")
                continue

            # Pure shape (no text) - only add if it has fill or stroke
            fill = _fill_from_shape(shape)
            stroke = _stroke_from_shape(shape)

            # Skip empty placeholders with no visual content
            if is_placeholder and not fill and not stroke:
                logger.info(f"  → Skipping empty placeholder shape {z} (no text frame)")
                continue

            if fill or stroke:
                shape_el = _shape_element(shape, z)
                elements.append(shape_el)
            else:
                logger.info(f"  → Skipping shape {z} with no fill or stroke")

        slides.append(
            ParsedSlide(
                width_px=slide_width_px,
                height_px=slide_height_px,
                background=background,
                fonts=list({f for f in fonts if f}),
                elements=elements,
                index=slide_index,
            )
        )

    return slides
