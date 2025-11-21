"""
Test script to debug fill extraction from PPTX shapes.
Run this to see what fill types your shapes have.
"""

import sys
from pptx import Presentation
from pptx.enum.dml import MSO_FILL

def test_fill_extraction(pptx_path: str):
    """Debug fill extraction for a PPTX file."""
    presentation = Presentation(pptx_path)

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        print(f"\n{'='*60}")
        print(f"SLIDE {slide_idx}")
        print(f"{'='*60}")

        for shape_idx, shape in enumerate(slide.shapes):
            print(f"\n  Shape {shape_idx}: {shape.shape_type}")

            # Check if shape has fill
            try:
                fill = shape.fill
                print(f"    ✓ Has fill object")

                try:
                    fill_type = fill.type
                    print(f"    Fill type: {fill_type}")

                    if fill_type == MSO_FILL.SOLID:
                        print(f"      → SOLID fill")
                        try:
                            color = fill.fore_color
                            rgb = color.rgb
                            hex_color = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                            print(f"      → Color: {hex_color}")
                        except Exception as e:
                            print(f"      → Could not extract color: {e}")

                    elif fill_type == MSO_FILL.BACKGROUND:
                        print(f"      → BACKGROUND fill (inherits from master)")
                        # Try to get the inherited color
                        try:
                            # Access the underlying XML to get theme color
                            print(f"      → This shape inherits background from slide master")
                        except Exception as e:
                            print(f"      → Error: {e}")

                    elif fill_type == MSO_FILL.GRADIENT:
                        print(f"      → GRADIENT fill")

                    elif fill_type == MSO_FILL.PATTERNED:
                        print(f"      → PATTERNED fill")

                    elif fill_type == MSO_FILL.PICTURE:
                        print(f"      → PICTURE fill")

                    else:
                        print(f"      → Unknown fill type: {fill_type}")

                except Exception as e:
                    print(f"    ✗ Could not get fill.type: {e}")

            except Exception as e:
                print(f"    ✗ No fill: {e}")

            # Check if shape has text
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                print(f"    ✓ Has text frame")
                try:
                    text = shape.text
                    if text:
                        print(f"    Text: {text[:50]}...")
                except:
                    pass

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python test_fill_extraction.py <path_to_pptx>")
        sys.exit(1)

    test_fill_extraction(sys.argv[1])
