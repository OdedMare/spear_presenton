from pptx import Presentation

def create_test_pptx(filename):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here"

    prs.save(filename)
    print(f"Saved {filename}")

if __name__ == "__main__":
    create_test_pptx("test_gen.pptx")
